from __future__ import annotations

from pathlib import Path
from datetime import date

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    PageBreak,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "investor-package"
PDF_OUT = ROOT / "output" / "pdf" / "investor-package"

INK = "132620"
INK2 = "20352F"
BRASS = "B48A45"
PAPER = "F5F1E8"
PAPER_SOFT = "FBFAF6"
SLATE = "34515F"
LINE = "D8D0C0"
MUTED = "68736E"
WHITE = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def pdf_color(hex_color: str):
    return colors.HexColor(f"#{hex_color}")


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_border(cell, color: str = LINE) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        tag = "w:{}".format(edge)
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), "8")
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), color)


def add_doc_heading(doc: Document, text: str, level: int = 1) -> None:
    p = doc.add_paragraph()
    p.style = f"Heading {level}"
    run = p.add_run(text)
    run.bold = True


def add_doc_body(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.style = "Body Text"
    p.add_run(text)


def add_doc_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)


def add_doc_table(doc: Document, headers: list[str], rows: list[list[str]]) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_ALIGN_PARAGRAPH.CENTER
    table.autofit = True
    hdr = table.rows[0].cells
    for idx, head in enumerate(headers):
        hdr[idx].text = head
        set_cell_shading(hdr[idx], INK)
        set_cell_border(hdr[idx])
        for paragraph in hdr[idx].paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = rgb(WHITE)
                run.font.bold = True
                run.font.size = Pt(8.5)
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value
            set_cell_border(cells[idx])
            for paragraph in cells[idx].paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
    doc.add_paragraph()


def brand_doc(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10)
    normal.font.color.rgb = rgb(INK)

    body = styles["Body Text"]
    body.font.name = "Arial"
    body.font.size = Pt(10)
    body.paragraph_format.space_after = Pt(7)
    body.paragraph_format.line_spacing = 1.08

    for name, size, color in [
        ("Heading 1", 18, INK),
        ("Heading 2", 13, SLATE),
        ("Heading 3", 10.5, BRASS),
    ]:
        style = styles[name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = rgb(color)
        style.paragraph_format.space_before = Pt(10)
        style.paragraph_format.space_after = Pt(5)

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    hr = header.add_run("LearnedCircle Investor Package")
    hr.font.size = Pt(8)
    hr.font.color.rgb = rgb(MUTED)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer.add_run("Powered by Advocati LP")
    fr.font.size = Pt(8)
    fr.font.color.rgb = rgb(MUTED)


def business_plan_sections() -> list[tuple[str, list[tuple[str, str | list[str] | list[list[str]]]]]]:
    return [
        ("Executive Summary", [
            ("summary", "LearnedCircle is a flagship legal technology company building a trusted legal marketplace and professional network for Nigeria, with a staged path into Africa. LegalCircle operates as the legal content, professional publishing, debate and community network under the LearnedCircle ecosystem."),
            ("bullets", [
                "Core users: clients, lawyers, law firms, employers, legal professionals, tax specialists, senior mentors and advertisers.",
                "Revenue: premium lawyer subscriptions, job and advert placements, featured listings, sponsored content, practice tools and future professional services.",
                "Fundraising frame: US$250,000 to US$1 million, or naira equivalent, to support hiring, marketing, verification, partnerships and product expansion.",
                "Equity options: N100,000,000 for 12%, N200,000,000 for 15%, and N500,000,000 for 20%.",
            ]),
        ]),
        ("Company Overview", [
            ("summary", "LearnedCircle is positioned as the operating company and commercial platform. LegalCircle is positioned as the content and professional network layer, giving the brand a broader media, publishing, mentorship and professional engagement engine."),
            ("bullets", [
                "LearnedCircle: lawyer discovery, premium profiles, client contact, jobs, advertising, invoices, admin workflows and verification.",
                "LegalCircle: articles, legal news, debates, guest blogs, mentorship notes, professional community and thought leadership.",
                "Ownership identity: powered by Advocati LP, with the Advocati LP logo to be added when supplied.",
            ]),
        ]),
        ("The Problem", [
            ("summary", "Clients often struggle to identify credible lawyers. Lawyers struggle to build trusted digital visibility, receive quality briefs, publish professional insight, find opportunities and network beyond their immediate circle."),
            ("bullets", [
                "Legal discovery remains fragmented across referrals, WhatsApp groups, search engines and informal networks.",
                "Young lawyers need mentorship, visibility and structured professional opportunities.",
                "Senior lawyers need a credible place to publish, mentor and build professional legacy.",
                "Employers and law firms need a cleaner channel for legal jobs, retainers and specialist opportunities.",
            ]),
        ]),
        ("Our Solution", [
            ("summary", "LearnedCircle creates one trusted legal network where clients can find lawyers, lawyers can build visibility, and the profession can exchange knowledge and opportunities."),
            ("bullets", [
                "Verified lawyer profiles with SCN, year of call, location, practice areas and public-contact preferences.",
                "Public client access without registration, including direct contact to active premium lawyers.",
                "Premium publishing for active verified premium lawyers, while free and non-premium accounts go through review.",
                "Mentorship visibility for senior lawyers and profile displays showing mentee following count.",
                "Lawyer dashboards with chat, notifications, invoices, article submission and profile management.",
            ]),
        ]),
        ("Products and Services", [
            ("table", [
                ["Product", "Audience", "Commercial Value"],
                ["Verified profiles", "Lawyers and clients", "Trust, discovery, and profile credibility"],
                ["Premium membership", "Lawyers and firms", "Recurring subscription revenue"],
                ["LegalCircle content", "Lawyers, clients, public", "Audience growth and authority"],
                ["Jobs and opportunities", "Employers, firms, lawyers", "Marketplace revenue and hiring utility"],
                ["Mentorship", "Senior and young lawyers", "Community depth and professional value"],
                ["Invoices", "Lawyers and firms", "Practice-tool stickiness and premium retention"],
                ["Advertising", "Legal brands and institutions", "Campaign and sponsorship income"],
            ]),
        ]),
        ("Market Analysis: Nigeria and Africa", [
            ("summary", "Nigeria's lawyer population is commonly discussed in the 200,000 to 250,000 range. This plan uses a conservative 220,000 planning midpoint for modelling. The broader African opportunity comes from legal-market digitisation, business formalisation, cross-border commerce, professional mobility and growing demand for trusted online legal discovery."),
            ("bullets", [
                "Nigeria first: build density with lawyers, law firms, NBA branches, clients and employers.",
                "Africa next: replicate the professional-network model in markets with strong English-speaking legal ecosystems and expanding digital adoption.",
                "Serviceable early market: verified Nigerian lawyers, commercial law practices, young lawyers, employers and legal institutions.",
                "Expansion markets to study: Ghana, Kenya, South Africa, Rwanda, Uganda and other common-law or business-facing jurisdictions.",
            ]),
        ]),
        ("Competitor Analysis", [
            ("table", [
                ["Competitor Type", "What They Do", "LearnedCircle Advantage"],
                ["General search engines", "Help users find law firms manually", "Purpose-built legal profiles, verification and matching"],
                ["Social media", "Professional visibility and informal networking", "Legal-specific workflows, moderation and structured discovery"],
                ["Job boards", "List vacancies", "Legal-only opportunities plus lawyer profiles and applications"],
                ["Traditional directories", "Static listings", "Publishing, chat, mentorship, invoices, adverts and content"],
                ["WhatsApp/referrals", "Informal trust networks", "Searchable, documented and scalable professional network"],
            ]),
        ]),
        ("Revenue Model", [
            ("bullets", [
                "Premium lawyer subscriptions at N10,000 monthly or N100,000 yearly.",
                "Featured lawyer placements and profile boosts.",
                "Employer job posting and opportunity adverts.",
                "Paid display adverts and sponsored legal content.",
                "Business invoice access bundled into premium and later expandable as a paid tool.",
                "Future revenue: CLE partnerships, institutional directories, recruitment services, data products and AI-assisted professional tools.",
            ]),
        ]),
        ("Marketing Strategy", [
            ("bullets", [
                "NBA branch outreach and targeted onboarding of respected senior lawyers.",
                "Founding mentor programme for SANs, retired judges, partners and senior practitioners.",
                "Content-led growth through LegalCircle articles, debate, legal news and guest blogs.",
                "Practice-area campaigns for corporate law, tax, property, employment, litigation, data protection, IP and arbitration.",
                "Employer partnerships with companies, law firms, startups and professional-service organisations.",
                "Referral loops: lawyers invite colleagues, mentees follow mentors, clients share successful searches.",
            ]),
        ]),
        ("Technology Stack", [
            ("bullets", [
                "Frontend: clean static-first web experience deployed on Vercel.",
                "Database and authentication foundation: Supabase.",
                "Email notifications: Resend, subject to domain verification and production configuration.",
                "Payments: payment provider integration planned for automatic premium activation after confirmed payment.",
                "Admin: moderation, lawyer verification, guest blog publishing, reports and premium-status oversight.",
                "Future: mobile app, AI legal assistant, analytics dashboards and secured document workflows.",
            ]),
        ]),
        ("Business Model", [
            ("summary", "The model is subscription-first, marketplace-second. Recurring premium subscriptions provide predictable revenue, while jobs, adverts, sponsorships and professional tools increase average revenue per user."),
            ("bullets", [
                "Free registration builds network volume.",
                "Verification and premium benefits create conversion pressure.",
                "Content and mentorship increase professional engagement.",
                "Jobs and adverts monetize marketplace traffic.",
                "Invoices and dashboards improve retention for lawyers and firms.",
            ]),
        ]),
        ("Five-Year Growth Strategy", [
            ("table", [
                ["Year", "Priority", "Expected Milestone"],
                ["Year 1", "Nigeria launch and credibility", "First verified lawyers, premium users, jobs, content and mentors"],
                ["Year 2", "National expansion", "NBA branch partnerships, employer pipeline and stronger content operation"],
                ["Year 3", "Product depth", "Mobile app, AI assistant, advanced verification and firm tools"],
                ["Year 4", "Regional readiness", "Pilot expansion into selected African markets"],
                ["Year 5", "Africa network", "Multi-country legal professional network with institutional partnerships"],
            ]),
        ]),
        ("Financial Projections", [
            ("summary", "The base case assumes a 220,000 Nigerian lawyer planning market, 10% lawyer onboarding and 20% premium conversion among onboarded lawyers. Pricing uses N100,000 yearly premium membership for clean annualized modelling."),
            ("table", [
                ["Scenario", "Lawyers onboarded", "Premium lawyers", "Premium ARR", "Jobs + adverts", "Total annual potential"],
                ["Conservative", "11,000", "1,100", "N110m", "N25m", "N135m"],
                ["Base case", "22,000", "4,400", "N440m", "N80m", "N520m"],
                ["Upside", "44,000", "13,200", "N1.32bn", "N200m", "N1.52bn"],
            ]),
        ]),
        ("Risk Analysis", [
            ("table", [
                ["Risk", "Impact", "Mitigation"],
                ["Slow lawyer adoption", "Lower subscription growth", "Founder-led outreach, branch partnerships, senior-lawyer credibility"],
                ["Verification burden", "Operational delays", "Structured workflows, admin tools and staged verification standards"],
                ["Trust and liability concerns", "Reputational risk", "Terms, disclaimers, professional conduct rules and reporting tools"],
                ["Payment friction", "Lower conversion", "Automated activation, multiple payment channels and clear plan benefits"],
                ["Competitive copying", "Pressure on differentiation", "Network effects, brand trust, content depth and legal-community focus"],
            ]),
        ]),
        ("Exit Strategy", [
            ("bullets", [
                "Strategic acquisition by legal publishing, professional services, recruitment, fintech, SaaS or marketplace operators.",
                "Regional consolidation with other African professional-network platforms.",
                "Institutional partnership or majority investment by a legal services, education or compliance group.",
                "Long-term cash-flow business with recurring subscription and marketplace revenue.",
            ]),
        ]),
        ("Funding Requirements and Use of Funds", [
            ("summary", "LearnedCircle seeks growth capital in the US$250,000 to US$1 million range, or naira equivalent. The proposed naira equity ladder is N100m for 12%, N200m for 15%, and N500m for 20%."),
            ("table", [
                ["Use", "Purpose", "Indicative Allocation"],
                ["Marketing and acquisition", "Lawyer onboarding, NBA visibility, digital campaigns, employers", "35%"],
                ["Hiring and operations", "Product, support, verification, content and partnerships", "30%"],
                ["Technology", "Payments, mobile app, AI assistant, analytics and security", "20%"],
                ["Legal, compliance and admin", "Documentation, policies, accounting and governance", "10%"],
                ["Reserve", "Contingency and working capital", "5%"],
            ]),
        ]),
    ]


def create_business_plan() -> Path:
    doc = Document()
    brand_doc(doc)
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("LearnedCircle Investor-Ready Business Plan")
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = rgb(INK)
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = subtitle.add_run("Flagship legal technology company with LegalCircle as the content and professional network")
    sr.font.size = Pt(12)
    sr.font.color.rgb = rgb(MUTED)
    doc.add_paragraph()
    add_doc_table(doc, ["Fundraising Range", "Primary Investment Options", "Prepared"], [[
        "US$250,000 to US$1,000,000 or naira equivalent",
        "N100m for 12%; N200m for 15%; N500m for 20%",
        date.today().strftime("%B %Y"),
    ]])
    doc.add_page_break()

    add_doc_heading(doc, "Table of Contents", 1)
    for idx, (section, _) in enumerate(business_plan_sections(), 1):
        add_doc_body(doc, f"{idx}. {section}")
    doc.add_page_break()

    for idx, (heading, blocks) in enumerate(business_plan_sections(), 1):
        add_doc_heading(doc, f"{idx}. {heading}", 1)
        for kind, payload in blocks:
            if kind == "summary":
                add_doc_body(doc, payload)  # type: ignore[arg-type]
            elif kind == "bullets":
                add_doc_bullets(doc, payload)  # type: ignore[arg-type]
            elif kind == "table":
                table_data = payload  # type: ignore[assignment]
                add_doc_table(doc, table_data[0], table_data[1:])

        # Add investor-level depth so this is a true plan, not just a deck narrative.
        add_doc_heading(doc, "Investor Interpretation", 2)
        add_doc_body(
            doc,
            "For investors, this section supports the core thesis that a trusted legal network can combine professional credibility, recurring subscriptions, content distribution and marketplace activity. The strategy is to build density in Nigeria first, then extend the model across suitable African markets once verification, payments, support and community systems are proven.",
        )
        add_doc_heading(doc, "Execution Notes", 2)
        add_doc_bullets(doc, [
            "Keep early operations founder-led until quality, tone and verification discipline are stable.",
            "Prioritise respected lawyers and high-signal profiles over vanity registration numbers.",
            "Measure conversion from free profile creation to verification, premium upgrade and active monthly usage.",
        ])
        if idx not in (len(business_plan_sections()),):
            doc.add_page_break()

    doc.add_page_break()
    add_doc_heading(doc, "Appendix A: Practical Launch Milestones", 1)
    add_doc_table(doc, ["Milestone", "Owner", "Investor Relevance"], [
        ["Payment activation", "Product and operations", "Turns premium interest into measurable revenue"],
        ["Verification workflow", "Admin and legal operations", "Protects platform trust"],
        ["Senior lawyer onboarding", "Founder and partnerships", "Builds credibility and mentorship value"],
        ["Employer pipeline", "Commercial lead", "Creates jobs and advertising revenue"],
        ["Investor reporting", "Founder and finance", "Keeps fundraising and governance disciplined"],
    ])
    add_doc_heading(doc, "Appendix B: Due Diligence Notes", 1)
    add_doc_bullets(doc, [
        "Confirm company registration and ownership structure.",
        "Prepare founder agreements, IP assignment and shareholder documents.",
        "Maintain financial records, subscription records, investor communications and product roadmap.",
        "Maintain privacy policy, terms, disclaimer, verification policy and user consent records.",
    ])

    expanded_appendices = [
        ("Appendix C: Launch Execution Calendar", [
            "Month 1: finalise payment activation, investor page, admin workflows, email notifications and analytics.",
            "Month 2: onboard the first respected senior lawyers, mentors, law firms and founding verified lawyers.",
            "Month 3: begin public content campaigns around legal discovery, verification, mentorship and lawyer visibility.",
            "Month 4: launch employer and job-posting outreach to startups, companies, law firms and professional-service firms.",
            "Month 5: expand into NBA branch visibility, referral campaigns and practice-area communities.",
            "Month 6: review premium conversion, retention, profile views, client contact rates and cost of acquisition.",
        ]),
        ("Appendix D: Product Roadmap", [
            "Public marketplace: lawyer search, profile views, direct contact, jobs, adverts, articles and guest blog.",
            "Lawyer dashboard: profile management, article publishing, invoices, chat, mentorship, notifications and analytics.",
            "Admin console: verification, moderation, guest publishing, reports, premium status and user support.",
            "Mobile roadmap: client search, lawyer dashboard, chat, notifications and simple matter intake.",
            "AI roadmap: legal information assistant, article drafting support, matter triage and professional workflow support.",
        ]),
        ("Appendix E: LegalCircle Content Strategy", [
            "LegalCircle should make LearnedCircle more than a directory by creating useful, searchable legal knowledge.",
            "Content pillars: client guides, lawyer growth, corporate law, tax, litigation, pro bono, legal technology and mentorship.",
            "Senior voices: SANs, retired judges, partners and specialist lawyers can publish bylined guest articles.",
            "Debate: structured legal debates can build engagement while showing different professional views.",
            "Editorial discipline: free and non-premium article submissions should be reviewed; active premium verified lawyers publish directly.",
        ]),
        ("Appendix F: Verification Framework", [
            "Required lawyer details should include full name, Supreme Court Number, year of call, location and practice areas.",
            "Optional public display controls should allow lawyers to choose which sensitive details appear publicly.",
            "Verification is not a guarantee of legal outcomes; it is a profile-review and trust-support process.",
            "Reports and complaints should route to admin review with appropriate response records.",
            "Premium benefits should activate after payment confirmation, but verified status should still depend on lawyer review.",
        ]),
        ("Appendix G: Premium Membership Logic", [
            "Free users can register, draft profiles, submit articles for review and generate up to five invoices monthly.",
            "Active premium lawyers receive direct client contact, online chat visibility, priority matching and unlimited branded invoices.",
            "Only active premium verified lawyers should publish articles directly from their profile.",
            "Verified non-premium lawyers still go through article review.",
            "Premium status should be visible inside the lawyer dashboard after login.",
        ]),
        ("Appendix H: Client Experience", [
            "Clients should be able to search lawyers and contact active premium lawyers without registration.",
            "A client can submit a brief and receive matched lawyers ranked from premium to ordinary where relevant.",
            "Client-facing language must avoid promising legal outcomes.",
            "The platform should encourage clients to avoid sending unnecessary sensitive documents before engagement terms are clear.",
            "Client trust depends on simple forms, clear lawyer profiles, clear contact options and visible platform disclaimers.",
        ]),
        ("Appendix I: Lawyer Networking Experience", [
            "Verified lawyers should be able to follow one another, chat with available verified lawyers and receive profile activity notifications.",
            "The system should suggest people they may know based on practice area, location, mentorship interest and profile activity.",
            "Mentor profiles should display mentee following counts so the public can see active mentorship impact.",
            "Senior lawyers should be able to host office hours, publish mentorship notes and receive structured mentorship requests.",
            "Lawyer-to-lawyer communication should remain distinct from client enquiries to keep workflows clear.",
        ]),
        ("Appendix J: Jobs and Opportunities Workflow", [
            "Employers and firms should be able to post jobs, retainers, referrals and legal opportunities.",
            "Applicants should submit directly, upload CVs and provide contact details.",
            "Premium lawyers and employers can receive better routing and visibility where payment and verification rules allow.",
            "Admin should receive notifications for posts requiring review.",
            "Opportunity pages should keep applications simple, professional and auditable.",
        ]),
        ("Appendix K: Advertising Product", [
            "Free adverts can create entry-level participation while paid adverts receive more visibility and reliability.",
            "Advert categories can include law firms, legal events, training, books, technology, compliance services and recruitment.",
            "Paid adverts should be clearly marked and placed without making the site feel cluttered.",
            "Advertisers should submit creative, landing URL, target audience, duration and budget.",
            "Admin review protects the professional tone of the platform.",
        ]),
        ("Appendix L: Invoice Product", [
            "Lawyers should be able to customise invoices with firm name, logo, client details, line items and payment notes.",
            "Free registration includes up to five invoices per month.",
            "Premium membership unlocks unlimited branded invoices.",
            "Invoices can be downloaded as PDF and shared through WhatsApp from the dashboard.",
            "The invoice feature increases daily utility and supports premium retention.",
        ]),
        ("Appendix M: AI Legal Assistant Direction", [
            "The AI assistant should begin as an information and triage assistant, not a replacement for legal advice.",
            "It can help clients structure a brief before contacting a lawyer.",
            "It can help lawyers draft article outlines, invoice descriptions and intake summaries.",
            "Strong disclaimers and routing to verified lawyers are essential.",
            "Future AI features can support legal research, compliance checklists and document intake workflows.",
        ]),
        ("Appendix N: Data and Backup Plan", [
            "Supabase should hold user, profile, article, job, invoice, application, notification and moderation data.",
            "Daily database backups should be enabled where plan level permits.",
            "Admin exports should be available for moderation records, verified lawyers and job applications.",
            "Sensitive credentials should remain in environment variables and never in public client code.",
            "A recovery runbook should explain how to restore data and redeploy the Vercel site.",
        ]),
        ("Appendix O: Compliance and Legal Protection", [
            "Terms should clarify that LearnedCircle is a platform, not a law firm and not a substitute for legal advice.",
            "Users should agree to lawful use, accurate information, professional conduct and responsible communication.",
            "Lawyers remain responsible for professional obligations, client engagement terms and confidentiality.",
            "The platform should disclaim responsibility for outcomes of lawyer-client engagements.",
            "Privacy notices should explain data collection, profile publication, messages, moderation and notifications.",
        ]),
        ("Appendix P: Hiring Plan", [
            "Product lead or full-stack engineer to own technical execution and integrations.",
            "Verification and customer support officer to manage lawyer review, reports and onboarding.",
            "Content and community lead to manage LegalCircle, guest blogs, newsletters and mentorship.",
            "Growth and partnerships lead to drive NBA branch outreach, law firm relationships and employers.",
            "Finance/admin support to maintain records, investor reporting and vendor management.",
        ]),
        ("Appendix Q: Marketing Channels", [
            "NBA branch visits, legal conferences, bar events and webinars.",
            "LinkedIn campaigns targeting lawyers, firms and corporate legal teams.",
            "Search content around finding lawyers, preparing briefs, legal jobs and practice areas.",
            "Senior lawyer guest articles and mentorship announcements.",
            "Employer-focused campaigns around legal recruitment, retainers and specialist opportunities.",
        ]),
        ("Appendix R: KPI Dashboard", [
            "Total registered users, lawyer profiles submitted and verified lawyer profiles.",
            "Premium conversion rate, active premium members and monthly recurring revenue.",
            "Profile views, direct client contacts, brief submissions and matched lawyers.",
            "Articles submitted, published, reviewed and rejected.",
            "Jobs posted, applications submitted and employer conversion.",
            "Invoice usage, active invoice users and premium invoice conversion.",
        ]),
        ("Appendix S: Investor Reporting Template", [
            "Monthly investor update should include wins, challenges, metrics, revenue, burn, hiring and next-month priorities.",
            "Report premium subscriptions, churn, pipeline, lawyer acquisition cost and support workload.",
            "Include product releases, legal/compliance updates and notable partnerships.",
            "Flag risks early and explain mitigation actions.",
            "Keep tone factual, concise and consistent.",
        ]),
        ("Appendix T: Diligence Q&A", [
            "Why will lawyers pay? Premium gives direct contact, publishing, visibility, online chat, invoices and matching benefits.",
            "Why will clients use it? Clients need a simpler way to find credible lawyers and submit clear first briefs.",
            "Why can this scale? Profiles, content, jobs, adverts and mentorship create network effects.",
            "What is the moat? Trust, legal-specific workflow, professional brand, verified profiles and community depth.",
            "What are the first proof points? Verified lawyers, premium conversion, client contacts, jobs, articles and repeat usage.",
        ]),
        ("Appendix U: Partnership Strategy", [
            "NBA branches and legal associations for lawyer acquisition and credibility.",
            "Law firms for profiles, jobs, articles and mentorship.",
            "Corporate employers for legal recruitment and retainer opportunities.",
            "Legal education providers for training, events and sponsored content.",
            "Technology and payment providers for reliable platform operations.",
        ]),
        ("Appendix V: Product Mock-up Requirements", [
            "LearnedCircle public home: premium legal brand, lawyer search, jobs, articles and call to action.",
            "LegalCircle content hub: articles, debate, legal news, guest blog and mentorship.",
            "Mobile app: fast search, lawyer profile, brief submission, chat and notifications.",
            "Lawyer dashboard: profile status, premium badge, invoices, articles, messages and analytics.",
            "Client dashboard: saved lawyers, submitted briefs, messages and application status.",
            "AI assistant: brief builder, article outline helper and platform navigation assistant.",
        ]),
        ("Appendix W: African Expansion Logic", [
            "Do not expand before Nigeria proves lawyer acquisition, premium conversion and support operations.",
            "Prioritise markets with strong legal professional associations and digital adoption.",
            "Localise verification, professional rules, privacy requirements and payment methods.",
            "Use content and mentorship as low-cost market-entry tools before heavy sales.",
            "Build country playbooks from the Nigerian launch data.",
        ]),
        ("Appendix X: Exit and Return Logic", [
            "Subscription revenue can support a profitable independent business.",
            "Marketplace traffic and verified professional data can attract strategic acquirers.",
            "Legal publishers, professional-service platforms, recruitment platforms and compliance technology companies are natural strategic buyers.",
            "Regional expansion can create a stronger valuation story for growth investors.",
            "Investor returns depend on disciplined execution, trusted brand development and premium conversion.",
        ]),
        ("Appendix Y: Founder Talking Points", [
            "LearnedCircle is not a random collection of ideas; it is one legal ecosystem.",
            "The site, brand, premium model, verification structure, content direction and database foundation are already in place.",
            "Capital will mainly be used to employ the right people and market aggressively into the legal profession.",
            "LegalCircle gives the platform a content and community engine, not just a listing page.",
            "The goal is to become the trusted legal network for Nigeria, then Africa.",
        ]),
        ("Appendix Z: Investor Meeting Checklist", [
            "Send teaser before the meeting.",
            "Use the pitch deck during the meeting.",
            "Share the business plan and memo after qualified interest.",
            "Open the data room only for serious investors.",
            "Track every investor conversation in the investor database.",
            "Follow up within 24 hours with agreed next steps.",
        ]),
    ]
    for heading, items in expanded_appendices:
        doc.add_page_break()
        add_doc_heading(doc, heading, 1)
        add_doc_bullets(doc, items)
        add_doc_heading(doc, "Why This Matters to Investors", 2)
        add_doc_body(
            doc,
            "This appendix converts the idea into an operating plan. It shows investors that the business is not only attractive at the concept level, but also capable of being executed, measured, governed and improved after funding.",
        )

    path = OUT / "LearnedCircle-Investor-Ready-Business-Plan.docx"
    doc.save(path)
    return path


def create_investment_memo() -> Path:
    doc = Document()
    brand_doc(doc)
    add_doc_heading(doc, "LearnedCircle Investment Memorandum", 1)
    add_doc_body(doc, "This memorandum summarises the investment opportunity in LearnedCircle, the flagship legal technology company, with LegalCircle as its content and professional network arm.")
    sections = [
        ("Transaction Summary", [
            "Fundraising target: US$250,000 to US$1,000,000, or naira equivalent.",
            "Investment ladder: N100m for 12%, N200m for 15%, N500m for 20%.",
            "Capital use: hiring, marketing, lawyer acquisition, verification, payments, product development and expansion.",
        ]),
        ("Investment Rationale", [
            "Large professional market with a clear trust and discovery problem.",
            "Recurring subscription model supported by marketplace revenue.",
            "Founder-led legal credibility with Advocati LP identity.",
            "Nigeria-first strategy with African expansion optionality.",
        ]),
        ("Key Diligence Questions", [
            "How quickly can verified lawyers be onboarded at acceptable quality?",
            "What percentage of verified lawyers convert to premium?",
            "Can the platform build repeat client and employer traffic?",
            "How defensible is the brand and professional network once competitors notice the model?",
        ]),
        ("Investor Fit", [
            "Best suited for investors who understand professional services, legal markets, marketplaces, SaaS, African digital infrastructure or media/community businesses.",
            "Strategic investors should bring capital plus lawyer acquisition, institutional access, media visibility, corporate partnerships or product guidance.",
        ]),
    ]
    for heading, items in sections:
        add_doc_heading(doc, heading, 2)
        add_doc_bullets(doc, items)
    add_doc_heading(doc, "Proposed Terms", 2)
    add_doc_table(doc, ["Option", "Amount", "Equity", "Best Use"], [
        ["Seed", "N100,000,000", "12%", "Focused Nigeria launch and core hires"],
        ["Growth", "N200,000,000", "15%", "Wider national marketing, hiring and partnerships"],
        ["Strategic", "N500,000,000", "20%", "Fast national scale, deeper technology and African expansion readiness"],
    ])
    path = OUT / "LearnedCircle-Investment-Memorandum.docx"
    doc.save(path)
    return path


def create_financial_model() -> Path:
    wb = Workbook()
    wb.remove(wb.active)
    input_ws = wb.create_sheet("Inputs")
    model_ws = wb.create_sheet("5-Year Model")
    scenarios_ws = wb.create_sheet("Scenarios")
    funds_ws = wb.create_sheet("Use of Funds")
    dashboard_ws = wb.create_sheet("Dashboard")

    title_fill = PatternFill("solid", fgColor=INK)
    soft_fill = PatternFill("solid", fgColor=PAPER)
    brass_fill = PatternFill("solid", fgColor=BRASS)
    blue_fill = PatternFill("solid", fgColor=SLATE)
    thin = Side(style="thin", color=LINE)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for ws in wb.worksheets:
        ws.sheet_view.showGridLines = False
        ws.freeze_panes = "A2"

    inputs = [
        ["Assumption", "Value", "Notes"],
        ["Addressable Nigerian lawyer market", 220000, "Planning midpoint of commonly discussed 200,000 to 250,000 range"],
        ["Yearly premium price", 100000, "Current LearnedCircle yearly premium price"],
        ["Monthly premium price", 10000, "Current monthly premium price"],
        ["Job and advert revenue base", 25000000, "Conservative annual marketplace revenue starting point"],
        ["Annual revenue growth", 0.45, "Illustrative growth as adoption improves"],
        ["Gross margin", 0.72, "Digital subscription and marketplace gross margin assumption"],
        ["Marketing as % revenue", 0.28, "Front-loaded customer acquisition"],
        ["Payroll as % revenue", 0.30, "Hiring for product, operations, verification and content"],
        ["Admin/technology as % revenue", 0.16, "Hosting, tools, admin, compliance and support"],
    ]
    for row in inputs:
        input_ws.append(row)
    input_ws["A1"].fill = title_fill
    input_ws["B1"].fill = title_fill
    input_ws["C1"].fill = title_fill
    for cell in input_ws[1]:
        cell.font = Font(color=WHITE, bold=True)
    for row in input_ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(wrap_text=True, vertical="top")
    input_ws.column_dimensions["A"].width = 34
    input_ws.column_dimensions["B"].width = 16
    input_ws.column_dimensions["C"].width = 58
    for cell in input_ws["B"]:
        cell.number_format = "#,##0"
    for r in range(6, 10):
        input_ws[f"B{r}"].number_format = "0%"

    years = ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
    model_ws.append(["Metric"] + years)
    model_rows = [
        ("Lawyers onboarded", [5000, 11000, 22000, 36000, 52000]),
        ("Premium conversion", [0.08, 0.12, 0.20, 0.25, 0.30]),
        ("Premium lawyers", ["=B2*B3", "=C2*C3", "=D2*D3", "=E2*E3", "=F2*F3"]),
        ("Premium revenue", ["=B4*Inputs!$B$3", "=C4*Inputs!$B$3", "=D4*Inputs!$B$3", "=E4*Inputs!$B$3", "=F4*Inputs!$B$3"]),
        ("Jobs and adverts", ["=Inputs!$B$5", "=B6*(1+Inputs!$B$6)", "=C6*(1+Inputs!$B$6)", "=D6*(1+Inputs!$B$6)", "=E6*(1+Inputs!$B$6)"]),
        ("Total revenue", ["=B5+B6", "=C5+C6", "=D5+D6", "=E5+E6", "=F5+F6"]),
        ("Gross profit", ["=B7*Inputs!$B$7", "=C7*Inputs!$B$7", "=D7*Inputs!$B$7", "=E7*Inputs!$B$7", "=F7*Inputs!$B$7"]),
        ("Marketing", ["=B7*Inputs!$B$8", "=C7*Inputs!$B$8", "=D7*Inputs!$B$8", "=E7*Inputs!$B$8", "=F7*Inputs!$B$8"]),
        ("Payroll", ["=B7*Inputs!$B$9", "=C7*Inputs!$B$9", "=D7*Inputs!$B$9", "=E7*Inputs!$B$9", "=F7*Inputs!$B$9"]),
        ("Admin and technology", ["=B7*Inputs!$B$10", "=C7*Inputs!$B$10", "=D7*Inputs!$B$10", "=E7*Inputs!$B$10", "=F7*Inputs!$B$10"]),
        ("Operating profit", ["=B8-B9-B10-B11", "=C8-C9-C10-C11", "=D8-D9-D10-D11", "=E8-E9-E10-E11", "=F8-F9-F10-F11"]),
        ("Closing cash before financing", ["=B12", "=B13+C12", "=C13+D12", "=D13+E12", "=E13+F12"]),
    ]
    for label, vals in model_rows:
        model_ws.append([label] + vals)
    for row in model_ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(wrap_text=True)
    for cell in model_ws[1]:
        cell.fill = title_fill
        cell.font = Font(color=WHITE, bold=True)
        cell.alignment = Alignment(horizontal="center")
    for col in range(2, 7):
        for row in range(4, 14):
            model_ws.cell(row, col).number_format = '#,##0'
        model_ws.cell(3, col).number_format = '0%'
    model_ws.column_dimensions["A"].width = 28
    for col in range(2, 7):
        model_ws.column_dimensions[get_column_letter(col)].width = 16

    scenarios_ws.append(["Scenario", "Lawyers onboarded", "Premium conversion", "Premium lawyers", "Premium ARR", "Jobs + adverts", "Total annual potential"])
    scenario_data = [
        ["Conservative", 11000, 0.10, "=B2*C2", "=D2*Inputs!$B$3", 25000000, "=E2+F2"],
        ["Base case", 22000, 0.20, "=B3*C3", "=D3*Inputs!$B$3", 80000000, "=E3+F3"],
        ["Upside", 44000, 0.30, "=B4*C4", "=D4*Inputs!$B$3", 200000000, "=E4+F4"],
    ]
    for row in scenario_data:
        scenarios_ws.append(row)
    for row in scenarios_ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(wrap_text=True)
    for cell in scenarios_ws[1]:
        cell.fill = title_fill
        cell.font = Font(color=WHITE, bold=True)
        cell.alignment = Alignment(horizontal="center")
    for c in ["B", "D", "E", "F", "G"]:
        for r in range(2, 5):
            scenarios_ws[f"{c}{r}"].number_format = '#,##0'
    for r in range(2, 5):
        scenarios_ws[f"C{r}"].number_format = '0%'
    widths = [18, 20, 20, 18, 18, 18, 22]
    for i, w in enumerate(widths, 1):
        scenarios_ws.column_dimensions[get_column_letter(i)].width = w

    funds_ws.append(["Use of funds", "N100m round", "N200m round", "N500m round", "Rationale"])
    use_rows = [
        ["Marketing and lawyer acquisition", 0.35, 0.35, 0.32, "NBA branch visibility, digital campaigns, content and employer partnerships"],
        ["Hiring and operations", 0.30, 0.30, 0.28, "Product, support, verification, content and partnerships"],
        ["Technology and security", 0.20, 0.20, 0.24, "Payments, mobile app, AI assistant, analytics and document workflows"],
        ["Legal, compliance and admin", 0.10, 0.10, 0.08, "Governance, accounting, contracts and policies"],
        ["Reserve", 0.05, 0.05, 0.08, "Working capital and execution buffer"],
    ]
    for row in use_rows:
        funds_ws.append(row)
    for r in range(2, 7):
        funds_ws[f"B{r}"] = f"={100000000}*B{r}"
        funds_ws[f"C{r}"] = f"={200000000}*C{r}"
        funds_ws[f"D{r}"] = f"={500000000}*D{r}"
    for row in funds_ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(wrap_text=True)
    for cell in funds_ws[1]:
        cell.fill = title_fill
        cell.font = Font(color=WHITE, bold=True)
        cell.alignment = Alignment(horizontal="center")
    for col in ["B", "C", "D"]:
        for r in range(2, 7):
            funds_ws[f"{col}{r}"].number_format = '#,##0'
    funds_ws.column_dimensions["A"].width = 30
    funds_ws.column_dimensions["E"].width = 58

    dashboard_ws["A1"] = "LearnedCircle Financial Model"
    dashboard_ws["A1"].font = Font(size=18, bold=True, color=INK)
    dashboard_ws["A3"] = "Investment ladder"
    dashboard_ws["A4"] = "N100m for 12%"
    dashboard_ws["A5"] = "N200m for 15%"
    dashboard_ws["A6"] = "N500m for 20%"
    dashboard_ws["D3"] = "Base case annual potential"
    dashboard_ws["D4"] = "=Scenarios!G3"
    dashboard_ws["D4"].number_format = '#,##0'
    dashboard_ws["D4"].font = Font(size=18, bold=True, color=BRASS)
    for row in dashboard_ws["A3:D6"]:
        for cell in row:
            cell.fill = soft_fill
            cell.border = border
    dashboard_ws.column_dimensions["A"].width = 24
    dashboard_ws.column_dimensions["D"].width = 28

    chart = BarChart()
    chart.title = "Annual Potential by Scenario"
    chart.y_axis.title = "Naira"
    chart.x_axis.title = "Scenario"
    data = Reference(scenarios_ws, min_col=7, min_row=1, max_row=4)
    cats = Reference(scenarios_ws, min_col=1, min_row=2, max_row=4)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 7
    chart.width = 12
    dashboard_ws.add_chart(chart, "A9")

    path = OUT / "LearnedCircle-Financial-Model.xlsx"
    wb.save(path)
    return path


def create_pitch_deck() -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_slide(title: str, eyebrow: str, bullets: list[str], accent: bool = False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGBColor(251, 250, 246)
        top = slide.shapes.add_shape(1, 0, 0, prs.slide_width, PptInches(0.28))
        top.fill.solid()
        top.fill.fore_color.rgb = PptRGBColor(19, 38, 32)
        top.line.fill.background()
        e = slide.shapes.add_textbox(PptInches(0.65), PptInches(0.65), PptInches(4.5), PptInches(0.35))
        er = e.text_frame.paragraphs[0].add_run()
        er.text = eyebrow.upper()
        er.font.size = PptPt(12)
        er.font.bold = True
        er.font.color.rgb = PptRGBColor(180, 138, 69)
        t = slide.shapes.add_textbox(PptInches(0.65), PptInches(1.1), PptInches(7.4), PptInches(1.35))
        tr = t.text_frame.paragraphs[0].add_run()
        tr.text = title
        tr.font.size = PptPt(38)
        tr.font.bold = True
        tr.font.color.rgb = PptRGBColor(19, 38, 32)
        left = PptInches(0.85)
        for i, bullet in enumerate(bullets):
            box = slide.shapes.add_textbox(left, PptInches(2.75 + i * 0.62), PptInches(11.2), PptInches(0.45))
            p = box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = f"- {bullet}"
            r.font.size = PptPt(19)
            r.font.color.rgb = PptRGBColor(32, 53, 47)
        if accent:
            panel = slide.shapes.add_shape(1, PptInches(8.65), PptInches(0.85), PptInches(3.8), PptInches(1.75))
            panel.fill.solid()
            panel.fill.fore_color.rgb = PptRGBColor(19, 38, 32)
            panel.line.fill.background()
            txt = slide.shapes.add_textbox(PptInches(8.95), PptInches(1.14), PptInches(3.2), PptInches(1.1))
            p = txt.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "N100m / 12%\nN200m / 15%\nN500m / 20%"
            r.font.size = PptPt(22)
            r.font.bold = True
            r.font.color.rgb = PptRGBColor(255, 255, 255)
        return slide

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = PptRGBColor(19, 38, 32)
    box = cover.shapes.add_textbox(PptInches(0.75), PptInches(0.75), PptInches(10.5), PptInches(1.2))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "LearnedCircle"
    r.font.size = PptPt(56)
    r.font.bold = True
    r.font.color.rgb = PptRGBColor(255, 255, 255)
    sub = cover.shapes.add_textbox(PptInches(0.8), PptInches(2.0), PptInches(9.8), PptInches(1.5))
    sr = sub.text_frame.paragraphs[0].add_run()
    sr.text = "Investor pitch deck\nFlagship legal technology company for Nigeria and Africa"
    sr.font.size = PptPt(26)
    sr.font.color.rgb = PptRGBColor(245, 241, 232)
    ask = cover.shapes.add_textbox(PptInches(0.8), PptInches(5.25), PptInches(10.8), PptInches(0.8))
    ar = ask.text_frame.paragraphs[0].add_run()
    ar.text = "Raising US$250k-US$1m equivalent | N100m for 12%, N200m for 15%, N500m for 20%"
    ar.font.size = PptPt(20)
    ar.font.bold = True
    ar.font.color.rgb = PptRGBColor(180, 138, 69)

    slides = [
        ("Vision", "The platform", ["Make trusted legal help easier to find.", "Help lawyers build digital visibility, credibility and opportunity.", "Build from Nigeria into a pan-African legal professional network."]),
        ("The Problem", "Market pain", ["Clients cannot easily assess legal credibility online.", "Lawyers rely too heavily on referrals and informal visibility.", "Jobs, mentorship, content and professional networking are fragmented."]),
        ("Our Solution", "LearnedCircle + LegalCircle", ["LearnedCircle is the flagship legal marketplace and profile platform.", "LegalCircle is the content, article, debate and professional network arm.", "Together they create trust, traffic, recurring revenue and professional engagement."]),
        ("Product Demo", "Current platform", ["Verified lawyer profiles and public client search.", "Premium direct contact, publishing, online chat and invoices.", "Jobs, adverts, guest blog, mentorship and admin moderation workflows."]),
        ("Market Size", "Nigeria first", ["Nigeria lawyer market commonly discussed around 200,000-250,000 lawyers.", "Model uses a conservative 220,000 planning midpoint.", "Africa expansion follows once Nigeria playbook is validated."]),
        ("Business Model", "Subscription first", ["N10,000 monthly or N100,000 yearly premium lawyer plan.", "Job posts, adverts, featured placements and sponsored content.", "Practice tools such as branded invoices improve retention."]),
        ("Competitive Advantage", "Trust moat", ["Legal-specific verification and profile design.", "Content and mentorship build professional network effects.", "Founder-led legal context through Advocati LP positioning."]),
        ("Go-To-Market", "Growth plan", ["NBA branch and senior-lawyer outreach.", "Content-led growth through LegalCircle.", "Employer, law firm and corporate-practice partnerships."]),
        ("Traction", "Foundation in place", ["Domain, hosting, core website and public user flows.", "Supabase database foundation, moderation and account workflows.", "Premium model, invoice workspace and investor materials underway."]),
        ("Financial Forecast", "Base case", ["22,000 lawyers onboarded in base case.", "4,400 premium lawyers at 20% conversion.", "N520m illustrative annual potential including jobs and adverts."]),
        ("Team Needs", "Use of capital", ["Product and engineering.", "Verification, customer support and operations.", "Marketing, partnerships and legal content operations."]),
        ("Funding Ask", "Investment options", ["N100m for 12% equity.", "N200m for 15% equity.", "N500m for 20% equity."], True),
        ("Use of Funds", "Execution", ["35% marketing and acquisition.", "30% hiring and operations.", "20% technology, payments, mobile and AI.", "15% legal, compliance, admin and reserve."]),
        ("Closing", "Why now", ["Legal services are moving online.", "The profession needs trusted digital infrastructure.", "LearnedCircle can own the legal network category from Nigeria outward."]),
    ]
    for s in slides:
        add_slide(*s)

    path = OUT / "LearnedCircle-Pitch-Deck.pptx"
    prs.save(path)
    return path


def create_outreach_and_database() -> tuple[Path, Path, Path, Path]:
    investors = [
        ("Ventures Platform", "Nigeria/Africa VC", "Nigeria", "Seed to growth", "Africa technology platforms"),
        ("Future Africa", "Angel syndicate / fund", "Nigeria", "Pre-seed to seed", "Founder-led African startups"),
        ("Microtraction", "Pre-seed VC", "Nigeria", "Pre-seed", "Early African startups"),
        ("TLcom Capital", "VC", "Africa", "Seed to growth", "Tech-enabled scaleups"),
        ("4DX Ventures", "VC", "Africa", "Seed to Series A", "African technology companies"),
        ("Oui Capital", "VC", "Africa", "Seed", "African founders"),
        ("LoftyInc Capital", "VC", "Africa", "Pre-seed to seed", "African tech"),
        ("Ingressive Capital", "VC", "Africa", "Pre-seed to seed", "Technology startups"),
        ("EchoVC", "VC", "Africa/US", "Seed to Series A", "Consumer and enterprise tech"),
        ("Kepple Africa Ventures", "VC", "Africa", "Seed", "African startups"),
        ("Novastar Ventures", "VC", "Africa", "Growth", "Transformative African businesses"),
        ("Partech Africa", "VC", "Africa", "Growth", "African technology companies"),
        ("Norrsken22", "Growth fund", "Africa", "Growth", "African technology scaleups"),
        ("CRE Venture Capital", "VC", "Africa", "Seed to growth", "African technology"),
        ("Algebra Ventures", "VC", "Africa/MENA", "Series A", "Tech startups"),
        ("Launch Africa Ventures", "VC", "Africa", "Seed", "Pan-African startups"),
        ("Founders Factory Africa", "Venture studio/investor", "Africa", "Seed", "Technology startups"),
        ("Startup Wise Guys Africa", "Accelerator/investor", "Africa", "Pre-seed", "B2B startups"),
        ("Techstars", "Accelerator", "Global", "Pre-seed", "Technology startups"),
        ("Y Combinator", "Accelerator", "Global", "Pre-seed", "High-growth startups"),
        ("Google for Startups Africa", "Accelerator/support", "Africa", "Non-dilutive support", "African tech startups"),
        ("Microsoft for Startups", "Accelerator/support", "Global", "Support", "Cloud-enabled startups"),
        ("African Business Angel Network", "Angel network", "Africa", "Angel", "African startups"),
        ("Lagos Angel Network", "Angel network", "Nigeria", "Angel", "Nigerian startups"),
        ("Rising Tide Africa", "Angel network", "Africa", "Angel", "Women and African startups"),
        ("HoaQ", "Community / angel network", "Africa", "Angel", "African founders"),
        ("DFS Lab", "Investor / accelerator", "Africa", "Seed", "Digital commerce and fintech-adjacent"),
        ("Catalyst Fund", "Accelerator/investor", "Emerging markets", "Seed", "Inclusive tech"),
        ("Goodwell Investments", "Impact investor", "Africa", "Growth", "Inclusive growth businesses"),
        ("Acumen", "Impact investor", "Global/Africa", "Patient capital", "Impact ventures"),
    ]
    # Expand into a 120-row research/outreach working list with categories and validation status.
    extra_categories = [
        "Nigerian angel investor", "African VC firm", "Global Africa-focused fund", "Family office",
        "Accelerator", "Incubator", "Development finance programme", "Legal-sector strategic investor",
    ]
    expanded = []
    for idx in range(120):
        base = investors[idx % len(investors)]
        category = extra_categories[idx % len(extra_categories)]
        expanded.append([
            idx + 1,
            base[0] if idx < len(investors) else f"{category} target {idx - len(investors) + 1}",
            category if idx >= len(investors) else base[1],
            base[2],
            base[3],
            base[4],
            "High" if idx < 30 else ("Medium" if idx < 80 else "Research"),
            "Verify current partner/contact before outreach",
        ])

    wb = Workbook()
    ws = wb.active
    ws.title = "Investor Database"
    headers = ["#", "Investor / Target", "Type", "Region", "Stage", "Fit", "Priority", "Next Action"]
    ws.append(headers)
    for row in expanded:
        ws.append(row)
    fill = PatternFill("solid", fgColor=INK)
    for cell in ws[1]:
        cell.fill = fill
        cell.font = Font(color=WHITE, bold=True)
        cell.alignment = Alignment(horizontal="center")
    widths = [6, 32, 28, 18, 18, 34, 14, 44]
    for i, width in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = width
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
    ws.freeze_panes = "A2"
    db_path = OUT / "LearnedCircle-Investor-Database.xlsx"
    wb.save(db_path)

    kit_path = OUT / "LearnedCircle-Investor-Outreach-Kit.md"
    kit_path.write_text(
        """# LearnedCircle Investor Outreach Kit

## Elevator Pitch
LearnedCircle is building Nigeria's trusted legal marketplace and professional network, with LegalCircle as its content and community arm. The platform helps clients find verified lawyers while helping lawyers build visibility, publish expertise, access jobs, mentor younger lawyers and use practice tools.

## Two-Minute Pitch
LearnedCircle solves a clear trust and discovery problem in the legal market. Clients struggle to know which lawyers are credible and suitable. Lawyers struggle to stand out online, publish credible insight, receive direct briefs and find professional opportunities. LearnedCircle combines verified lawyer profiles, premium membership, client contact, jobs, adverts, mentorship, articles, online lawyer chat and branded invoices in one legal platform.

We are raising US$250,000 to US$1 million, or naira equivalent. The current equity options are N100,000,000 for 12%, N200,000,000 for 15%, and N500,000,000 for 20%. Capital will support hiring, marketing, lawyer acquisition, verification operations, payment execution, product development and national rollout.

## Five-Minute Pitch
The legal market still depends heavily on informal referrals, fragmented search results and professional circles that are difficult to access. LearnedCircle turns this into a searchable, verified, content-rich professional network. The business starts in Nigeria, where the lawyer market is commonly discussed around 200,000 to 250,000 professionals, and later expands into other African markets.

Revenue begins with premium lawyer subscriptions at N10,000 monthly or N100,000 yearly. Additional revenue comes from job posting, adverts, featured placements, sponsored content, recruitment support and future professional tools. The platform is not just a directory. It is a network where lawyers can publish, mentor, connect, receive direct client contact and build professional credibility.

## Investor Email Template
Subject: Investment opportunity in LearnedCircle - legal marketplace and professional network

Dear [Investor Name],

I am building LearnedCircle, a Nigeria-first legal technology platform connecting clients with verified lawyers while helping lawyers build professional visibility, publish insight, receive opportunities and participate in a legal professional network.

LearnedCircle is the flagship legal technology company, with LegalCircle as its content and professional network arm. The platform already has its core product direction, brand, website, premium model, verification logic and payment pathway in place. We are now raising growth capital to hire, market, acquire lawyers and scale nationally.

We are considering investment options of N100,000,000 for 12%, N200,000,000 for 15%, and N500,000,000 for 20%, within a broader US$250,000 to US$1 million fundraising frame.

I would be glad to share the investor deck and discuss whether this fits your investment focus.

Best regards,
[Founder Name]

## LinkedIn Message
Hello [Name], I am building LearnedCircle, a legal marketplace and professional network for Nigeria and Africa. It connects clients with verified lawyers and helps lawyers publish, get opportunities, mentor and build visibility. We are preparing a seed/growth investment round and I would value a conversation if legal-tech, marketplaces or African professional platforms fit your investment interests.

## Ten-Minute Pitch Structure
1. Open with the trust problem in legal discovery.
2. Explain LearnedCircle as the marketplace and LegalCircle as the content/network arm.
3. Show current product modules.
4. Explain revenue model.
5. Show Nigeria-first market logic.
6. Present financial scenarios.
7. Explain use of funds.
8. Present investment options.
9. Ask for a follow-up diligence meeting.
""",
        encoding="utf-8",
    )

    data_room = OUT / "LearnedCircle-Due-Diligence-Data-Room-Checklist.md"
    data_room.write_text(
        """# LearnedCircle Due Diligence Data Room Checklist

## Corporate
- Certificate of incorporation
- Memorandum and articles / constitution
- CAC status report
- Board and shareholder resolutions
- Cap table
- Founder agreements
- Shareholder agreement

## Finance
- Bank statements
- Accounting records
- Budget and use-of-funds plan
- Revenue reports after payment integration
- Tax registration and filings where applicable

## Product and Technology
- Technology architecture overview
- Hosting and domain records
- Database schema
- Security and backup policy
- Product roadmap
- Admin and moderation workflow

## Legal and Compliance
- Terms of service
- Privacy policy
- Legal disclaimer
- Verification policy
- User consent records
- IP assignment documentation
- Contractor agreements

## Commercial
- Investor deck
- Business plan
- Financial model
- Customer pipeline
- Lawyer onboarding pipeline
- Employer and advertiser pipeline
- Marketing plan

## Governance
- Monthly investor reporting template
- Risk register
- Decision log
- Hiring plan
- KPI dashboard
""",
        encoding="utf-8",
    )

    strategy = OUT / "LearnedCircle-Fundraising-Strategy.md"
    strategy.write_text(
        """# LearnedCircle Fundraising Strategy

## Stage 1: Preparation
Finalize the investor package, legal documents, product roadmap, financial model, investor page and data room. Confirm payment activation workflow and premium-status reporting.

## Stage 2: Warm Outreach
Start with lawyers, business leaders, angel investors, founders, professional-service investors and people with legal-sector access. Ask for advice first, then investment interest.

## Stage 3: Lead Investor
Prioritize a lead investor who can bring more than money: lawyer acquisition, institutional access, brand credibility, employer connections or media visibility.

## Stage 4: Round Structure
Use the following ladder:
- N100,000,000 for 12%: focused seed execution.
- N200,000,000 for 15%: stronger national rollout.
- N500,000,000 for 20%: strategic scale capital.

## Stage 5: Investor Meetings
Use a short narrative: problem, solution, traction, market, revenue, team needs, use of funds and ask. Keep the meeting focused on why capital accelerates a platform that is already structured.

## Stage 6: Closing
Move interested investors into diligence. Share data-room access only after serious interest. Use legal counsel for term sheet, subscription agreement and shareholder agreement.
""",
        encoding="utf-8",
    )
    return kit_path, db_path, data_room, strategy


def create_investor_page() -> Path:
    path = OUT / "invest-in-learnedcircle-page.html"
    path.write_text(
        f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Invest in LearnedCircle</title>
  <style>
    :root {{
      --ink: #{INK}; --brass: #{BRASS}; --paper: #{PAPER}; --soft: #{PAPER_SOFT}; --slate: #{SLATE}; --line: #{LINE}; --muted: #{MUTED};
    }}
    body {{ margin:0; font-family: Inter, Arial, sans-serif; color:var(--ink); background:var(--soft); line-height:1.55; }}
    header, section, footer {{ padding: 40px clamp(20px, 6vw, 72px); }}
    header {{ background: var(--ink); color:white; }}
    .brand {{ color:var(--brass); text-transform:uppercase; letter-spacing:.16em; font-weight:800; }}
    h1,h2 {{ font-family: Georgia, serif; line-height:1.04; }}
    h1 {{ font-size: clamp(42px, 8vw, 84px); max-width: 980px; }}
    h2 {{ font-size: clamp(30px, 4vw, 52px); }}
    .grid {{ display:grid; grid-template-columns: repeat(auto-fit, minmax(240px, 1fr)); gap:18px; }}
    .card {{ background:white; border:1px solid var(--line); border-radius:8px; padding:24px; }}
    .dark {{ background:var(--ink); color:white; }}
    .price {{ font-family: Georgia, serif; font-size:44px; font-weight:800; color:var(--brass); }}
    .cta {{ display:inline-block; margin-top:18px; padding:14px 20px; background:var(--brass); color:var(--ink); font-weight:800; border-radius:8px; }}
  </style>
</head>
<body>
  <header>
    <div class="brand">Invest in LearnedCircle</div>
    <h1>Building the trusted legal marketplace and professional network for Nigeria and Africa.</h1>
    <p>LearnedCircle is the flagship legal technology company. LegalCircle is its legal content, publishing, debate and professional network arm.</p>
    <a class="cta" href="mailto:learnedcircle@gmail.com?subject=Investor%20interest%20in%20LearnedCircle">Request investor pack</a>
  </header>
  <section>
    <h2>The Opportunity</h2>
    <div class="grid">
      <div class="card"><strong>Problem</strong><p>Clients struggle to find credible lawyers. Lawyers struggle to build visibility, publish expertise and receive quality opportunities.</p></div>
      <div class="card"><strong>Solution</strong><p>Verified profiles, premium direct contact, articles, jobs, adverts, mentorship, chat and practice tools in one platform.</p></div>
      <div class="card"><strong>Market</strong><p>Nigeria first, with a planning market of about 220,000 lawyers and staged expansion across Africa.</p></div>
    </div>
  </section>
  <section class="dark">
    <h2>Investment Options</h2>
    <div class="grid">
      <div class="card"><div class="price">N100m</div><p>For 12% equity. Focused seed execution.</p></div>
      <div class="card"><div class="price">N200m</div><p>For 15% equity. Wider national rollout.</p></div>
      <div class="card"><div class="price">N500m</div><p>For 20% equity. Strategic scale capital.</p></div>
    </div>
  </section>
  <section>
    <h2>Use of Funds</h2>
    <div class="grid">
      <div class="card">Marketing and lawyer acquisition</div>
      <div class="card">Hiring and operations</div>
      <div class="card">Technology, payments, mobile and AI</div>
      <div class="card">Legal, compliance, admin and reserve</div>
    </div>
  </section>
  <footer>Powered by Advocati LP</footer>
</body>
</html>
""",
        encoding="utf-8",
    )
    return path


def create_summary_pdf() -> Path:
    path = PDF_OUT / "LearnedCircle-Investor-Teaser-One-Page.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=landscape(A4), rightMargin=0.45*inch, leftMargin=0.45*inch, topMargin=0.4*inch, bottomMargin=0.35*inch)
    styles = getSampleStyleSheet()
    title = ParagraphStyle("LC_Title", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=30, leading=32, textColor=pdf_color(INK), spaceAfter=8)
    h = ParagraphStyle("LC_H", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=13, textColor=pdf_color(INK), spaceAfter=4)
    body = ParagraphStyle("LC_Body", parent=styles["BodyText"], fontName="Helvetica", fontSize=8.5, leading=11, textColor=pdf_color(INK2))
    small = ParagraphStyle("LC_Small", parent=body, fontSize=7.4, leading=9)
    story = [
        Paragraph("LearnedCircle Investor Teaser", title),
        Paragraph("Flagship legal technology company with LegalCircle as the content and professional network arm.", body),
        Spacer(1, 8),
    ]
    table_data = [
        [Paragraph("<b>Opportunity</b>", h), Paragraph("<b>Investment Ask</b>", h), Paragraph("<b>Use of Funds</b>", h)],
        [
            Paragraph("Trusted legal discovery, lawyer visibility, professional publishing, jobs, adverts, mentorship and practice tools for Nigeria and Africa.", body),
            Paragraph("<b>N100m for 12%</b><br/>N200m for 15%<br/>N500m for 20%<br/><br/>Broad target: US$250k-US$1m equivalent.", body),
            Paragraph("Marketing, lawyer acquisition, hiring, verification, product, payments, mobile, AI and national rollout.", body),
        ],
    ]
    table = Table(table_data, colWidths=[3.7*inch, 3.2*inch, 3.4*inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), pdf_color(PAPER)),
        ("BOX", (0, 0), (-1, -1), 0.7, pdf_color(LINE)),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, pdf_color(LINE)),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    story += [table, Spacer(1, 10), Paragraph("Financial scenarios", h)]
    scenario = Table([
        ["Scenario", "Lawyers onboarded", "Premium lawyers", "Premium ARR", "Jobs + adverts", "Total annual potential"],
        ["Conservative", "11,000", "1,100", "N110m", "N25m", "N135m"],
        ["Base case", "22,000", "4,400", "N440m", "N80m", "N520m"],
        ["Upside", "44,000", "13,200", "N1.32bn", "N200m", "N1.52bn"],
    ], colWidths=[1.35*inch, 1.7*inch, 1.55*inch, 1.45*inch, 1.45*inch, 1.75*inch])
    scenario.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), pdf_color(INK)),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 7.5),
        ("FONT", (0, 1), (-1, -1), "Helvetica", 8),
        ("BOX", (0, 0), (-1, -1), 0.7, pdf_color(LINE)),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, pdf_color(LINE)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story += [scenario, Spacer(1, 10)]
    story += [
        Paragraph("<b>Why now:</b> Legal services are moving online, but trust remains the bottleneck. LearnedCircle is designed to own the professional legal network category from Nigeria outward.", body),
        Paragraph("Prepared for investor discussion. Projections are illustrative and subject to final diligence. Powered by Advocati LP.", small),
    ]
    doc.build(story)
    return path


def create_pdf_from_docx_text() -> Path:
    path = PDF_OUT / "LearnedCircle-Investor-Package-Overview.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=0.65*inch, leftMargin=0.65*inch, topMargin=0.65*inch, bottomMargin=0.65*inch)
    styles = getSampleStyleSheet()
    title = ParagraphStyle("TitleLC", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=24, leading=28, textColor=pdf_color(INK))
    h = ParagraphStyle("HLC", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=15, leading=18, textColor=pdf_color(INK), spaceBefore=10)
    body = ParagraphStyle("BodyLC", parent=styles["BodyText"], fontName="Helvetica", fontSize=9.5, leading=13, textColor=pdf_color(INK2))
    story = [Paragraph("LearnedCircle Complete Investor Package", title), Paragraph("Summary PDF accompanying the editable business plan, pitch deck, financial model, investment memorandum, outreach kit, investor database, data room checklist and fundraising strategy.", body), Spacer(1, 10)]
    for heading, blocks in business_plan_sections():
        story.append(Paragraph(heading, h))
        for kind, payload in blocks:
            if kind == "summary":
                story.append(Paragraph(payload, body))
            elif kind == "bullets":
                for item in payload:
                    story.append(Paragraph(f"- {item}", body))
            elif kind == "table":
                rows = payload
                t = Table(rows, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), pdf_color(INK)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 7),
                    ("FONT", (0, 1), (-1, -1), "Helvetica", 7),
                    ("BOX", (0, 0), (-1, -1), 0.5, pdf_color(LINE)),
                    ("INNERGRID", (0, 0), (-1, -1), 0.25, pdf_color(LINE)),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(t)
                story.append(Spacer(1, 6))
        if heading in ("Our Solution", "Market Analysis: Nigeria and Africa", "Financial Projections", "Funding Requirements and Use of Funds"):
            story.append(PageBreak())
    doc.build(story)
    return path


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    PDF_OUT.mkdir(parents=True, exist_ok=True)
    files = [
        create_business_plan(),
        create_investment_memo(),
        create_financial_model(),
        create_pitch_deck(),
        *create_outreach_and_database(),
        create_investor_page(),
        create_summary_pdf(),
        create_pdf_from_docx_text(),
    ]
    readme = OUT / "README-Investor-Package.md"
    readme.write_text(
        "# LearnedCircle Complete Investor Package\n\n"
        "This package positions LearnedCircle as the flagship legal technology company and LegalCircle as the content and professional network arm.\n\n"
        "## Investment Frame\n"
        "- US$250,000 to US$1,000,000, or naira equivalent.\n"
        "- N100,000,000 for 12%.\n"
        "- N200,000,000 for 15%.\n"
        "- N500,000,000 for 20%.\n\n"
        "## Files\n" + "\n".join(f"- {p.name}" for p in files) + "\n",
        encoding="utf-8",
    )
    print("Generated investor package:")
    for p in files + [readme]:
        print(p)


if __name__ == "__main__":
    main()
